#!/usr/bin/env python3
"""Extract text and notes from presentations into Markdown.

Supports:
- .pptx via python-pptx
- .pdf via pdftotext

Usage:
    python3 g.Presentations/scripts/extract.py

Outputs Markdown files under g.Presentations/extracted/.
"""

import subprocess
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "curated"
OUT = ROOT / "extracted"
OUT.mkdir(exist_ok=True)


def extract_pptx(path: Path) -> str:
    prs = Presentation(path)
    lines = [f"# {path.stem}", "", f"Source: `{path}`", f"Slides: {len(prs.slides)}", ""]
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("")
                lines.append("**Notes:**")
                lines.append(notes)
        lines.append("")
    return "\n".join(lines)


def extract_pdf(path: Path) -> str:
    result = subprocess.run(
        ["pdftotext", "-layout", str(path), "-"],
        capture_output=True,
        text=True,
        check=True,
    )
    text = result.stdout.strip()
    lines = [
        f"# {path.stem}",
        "",
        f"Source: `{path}`",
        f"Format: PDF",
        "",
        "```",
        text,
        "```",
        "",
    ]
    return "\n".join(lines)


def main():
    files = sorted(SOURCE.iterdir())
    for path in files:
        if path.suffix.lower() == ".pptx":
            md = extract_pptx(path)
            out_path = OUT / f"{path.stem}.pptx.md"
        elif path.suffix.lower() == ".pdf":
            md = extract_pdf(path)
            out_path = OUT / f"{path.stem}.pdf.md"
        else:
            continue
        out_path.write_text(md, encoding="utf-8")
        print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
